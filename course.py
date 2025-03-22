import os
import google.generativeai as genai
import gradio as gr
import glob
import pandas as pd
from docx import Document
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_google_genai import ChatGoogleGenerativeAI

# **1️⃣ Set API Key (Ensure secrets.toml is configured for security)**
os.environ['GOOGLE_API_KEY'] = "AIzaSyBMx_ZelxjCy6zNnaaArj78xd1rx8VWTdA"

# **2️⃣ Define Folder Path to Read Files from GitHub Repository Root**
DOCUMENTS_FOLDER_PATH = os.getcwd()  # Uses the directory where this script is placed

# **3️⃣ Function to Extract Text from Different Document Types**
def extract_text_from_file(file_path):
    """Extracts text from PDF, DOCX, PPTX, and XLSX files."""
    text = ""

    if file_path.endswith(".pdf"):
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        text = "\n".join([doc.page_content for doc in documents])

    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])

    elif file_path.endswith(".ppt") or file_path.endswith(".pptx"):
        presentation = Presentation(file_path)
        text_list = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():  # Ensure shape has text
                    text_list.append(shape.text.strip())  # Append clean text
        text = "\n".join(text_list)

    elif file_path.endswith(".xls") or file_path.endswith(".xlsx"):
        df = pd.read_excel(file_path)
        text = df.to_string()  # Convert entire spreadsheet to string

    return text.strip()

# **4️⃣ Function to Load and Process All Supported Files from GitHub Root Folder**
def load_and_process_documents(folder_path):
    """Loads and processes all PDFs, DOCX, PPTX, and XLSX files from the folder."""
    all_texts = []
    supported_formats = ["pdf", "docx", "ppt", "pptx", "xls", "xlsx"]
    files = sum([glob.glob(f"{folder_path}/*.{ext}") for ext in supported_formats], [])

    if not files:
        raise ValueError("⚠️ No supported files found in the specified folder!")

    for file in files:
        extracted_text = extract_text_from_file(file)
        if extracted_text:
            all_texts.append(extracted_text)

    # **5️⃣ Chunking & Vectorization**
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    docs = text_splitter.create_documents(all_texts)

    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-V2")
    vectorstore = FAISS.from_documents(docs, embeddings)

    return vectorstore

# **6️⃣ Load and Process Documents at Startup**
vectorstore = load_and_process_documents(DOCUMENTS_FOLDER_PATH)

# **7️⃣ AI-Powered Question Answering**
def ask_question(query):
    if not vectorstore:
        return "⚠️ No documents available. Please check the folder and restart the application."

    llm = ChatGoogleGenerativeAI(model="gemini-2.0-flash-exp")
    qa = RetrievalQA.from_chain_type(llm, chain_type="stuff", retriever=vectorstore.as_retriever())

    answer = qa.run(query)
    return f"「 ✦ ⚙️AI ✦ 」: {answer}"

# **🔹 8️⃣ Gradio Interface**
with gr.Blocks() as demo:
    gr.Markdown("<h2 style='text-align: center; color: blue;'># 「 ✦ ⚙️AI-Powered Course Management ✦ 」</h2>")
    gr.Markdown("This system is preloaded with PDF, Word, PowerPoint, and Excel files from the GitHub root folder. Ask any question based on them.")

    query_input = gr.Textbox(label="Ask a question")
    submit_button = gr.Button("💬 Get Answer")
    output_text = gr.Textbox(label="AI Response", interactive=False)

    submit_button.click(ask_question, inputs=[query_input], outputs=[output_text])

# **9️⃣ Launch the Gradio Interface**
demo.launch()